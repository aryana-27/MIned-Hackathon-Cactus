import streamlit as st
import tempfile, os, time, re, textwrap, requests
import openai
import google.generativeai as genai
from gtts import gTTS
from PyPDF2 import PdfReader
from pptx import Presentation
from moviepy.editor import VideoFileClip, concatenate_videoclips, TextClip, AudioFileClip
from moviepy import editor


# -------------------------
# PAGE CONFIGURATION
# -------------------------
st.set_page_config(page_title="Research Paper Pipeline", layout="wide")

# -------------------------
# API KEY SETUP
# -------------------------
# OpenAI API Key
openai.api_key = "YOUR OPENAI API KEY"
if not openai.api_key:
    st.error("Missing OpenAI API Key! Please set the OPENAI_API_KEY environment variable.")

# Gemini API Key
os.environ["GEMINI_API_KEY"] = os.getenv("GEMINI_API_KEY", "YOUR GEMINI API KEY")  # Replace with your key if needed
GEN_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEN_API_KEY:
    st.error("Missing Gemini API Key! Please set the GEMINI_API_KEY environment variable.")
genai.configure(api_key=GEN_API_KEY)

# Optionally, set your Pexels API Key as an environment variable named "PEXELS_API_KEY"

# -------------------------
# PDF TEXT EXTRACTION
# -------------------------
def extract_text_without_references(file_path):
    """
    Extract text from a PDF file, stopping before a References/Bibliography section.
    """
    try:
        reader = PdfReader(file_path)
        extracted_text = ""
        stop_keywords = ["References", "Bibliography", "Citations"]
        found_references = False
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                # If a stop keyword is found on the page, stop adding further text.
                if any(keyword.lower() in page_text.lower() for keyword in stop_keywords):
                    found_references = True
                elif not found_references:
                    extracted_text += page_text + "\n\n"
        return extracted_text
    except Exception as e:
        st.error(f"Error extracting text: {e}")
        return ""

# -------------------------
# UTILITY FUNCTIONS
# -------------------------
def split_text(text, max_chunk_size=1000):
    """Split text into chunks up to max_chunk_size characters."""
    return textwrap.wrap(text, width=max_chunk_size, break_long_words=False, replace_whitespace=False)

def extract_references(text):
    """Extract references (e.g., Table 1, Figure 2) from the text."""
    ref_pattern = r"(Table|Figure|Equation|Fig|Eq)\s?\d+"
    refs = re.findall(ref_pattern, text)
    return list(set(refs))

# -------------------------
# BULLET POINT SUMMARIZATION PIPELINE (using OpenAI)
# -------------------------
def summarize_text(text):
    """
    Generate a bullet-point summary from the provided research text using the OpenAI API.
    """
    prompt = f"""
You are an expert academic summarizer. Your task is to extract the key points from the research text provided below. 
Please generate a bullet point summary where each bullet point is clear, concise, and focused on a distinct idea or result.
Ensure that the summary:
- Contains only bullet points (each starting with a dash '-' or similar).
- Is written in a formal yet accessible tone.
- Covers all important aspects without repetition.
If applicable, provide at least five bullet points.

Research Text:
{text}

Bullet Points:
"""
    try:
        response = openai.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant specialized in academic summarization."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.5,
        )
        summary = response.choices[0].message['content'].strip()
        return summary
    except Exception as e:
        st.error(f"Error during summarization: {e}")
        return ""

def summarize_research_paper(text):
    """
    Process the research paper text by splitting it into manageable chunks,
    generating bullet-point summaries for each chunk, and collating references.
    """
    chunks = split_text(text, max_chunk_size=800)
    full_summary = ""
    all_refs = set()
    for i, chunk in enumerate(chunks):
        st.write(f"Processing Chunk {i+1}/{len(chunks)}...")
        refs = extract_references(chunk)
        summary = summarize_text(chunk)
        st.write(summary)
        full_summary += "\n" + summary
        all_refs.update(refs)
    return full_summary, all_refs

# -------------------------
# PPT GENERATION PIPELINE
# -------------------------
def create_ppt_from_text(text):
    """
    Create a PowerPoint presentation from the provided text.
    The text is split by double newlines, with the first line as the title.
    """
    ppt = Presentation()
    slides = text.strip().split('\n\n')
    for slide_text in slides:
        slide_layout = ppt.slide_layouts[1]  # Title & Content layout
        slide = ppt.slides.add_slide(slide_layout)
        lines = slide_text.split('\n')
        slide.shapes.title.text = lines[0] if lines else "Slide Title"
        slide.placeholders[1].text = "\n".join(lines[1:]) if len(lines) > 1 else ""
    ppt_path = "output.pptx"
    ppt.save(ppt_path)
    return ppt_path

# -------------------------
# PODCAST PRODUCTION PIPELINE (using gTTS for Audio)
# -------------------------
def split_text_for_podcast(text, max_chunk_size=800):
    return textwrap.wrap(text, width=max_chunk_size, break_long_words=False, replace_whitespace=False)

def generate_podcast_script(text):
    """
    Use the Gemini API to generate a 3-minute podcast script from research text.
    The script should be natural, conversational, and include intonation cues.
    """
    try:
        model = genai.GenerativeModel("gemini-pro")
        prompt = f"""
You are an expert monologue podcast scriptwriter who creates AI-friendly scripts.
Generate a natural-sounding 3-minute podcast script (~450-600 words) with a conversational tone.
Avoid headings like 'Introduction' or 'Conclusion'. Use pauses (… or [pause]) and intonation cues.
Use rhetorical questions, anecdotes, and friendly narration.
Research Text: {text}
"""
        response = model.generate_content(prompt)
        if response and response.text:
            return response.text
        else:
            return None
    except Exception as e:
        st.error(f"Error generating podcast script: {e}")
        return None

def process_research_for_podcast(text):
    """
    Process the research text into chunks, generate a podcast script for each,
    and collate the results.
    """
    chunks = split_text_for_podcast(text, max_chunk_size=800)
    full_script = ""
    all_refs = set()
    for i, chunk in enumerate(chunks):
        st.write(f"Processing Podcast Chunk {i+1}/{len(chunks)}...")
        refs = extract_references(chunk)
        script = generate_podcast_script(chunk)
        if script:
            st.write(script)
            full_script += "\n" + script
        all_refs.update(refs)
    return full_script, all_refs

def generate_gtts_audio(text, output_path="podcast_audio.mp3", lang="en"):
    """
    Generate an MP3 audio file from text using gTTS.
    """
    try:
        tts = gTTS(text=text, lang=lang)
        tts.save(output_path)
        st.write(f"Audio saved as {output_path}")
        return output_path
    except Exception as e:
        st.error(f"Error generating gTTS audio: {e}")
        return None

# -------------------------
# VIDEO CREATION PIPELINE (Simplified Example)
# -------------------------
def generate_1min_storyboard(text):
    """
    Use the Gemini API to generate a 1-minute storyboard for an Instagram Reel.
    """
    try:
        model = genai.GenerativeModel("gemini-pro")
        prompt = f"""
You are an expert in video storytelling. Generate a structured cinematic storyboard for a 1-minute Instagram Reel (60 seconds total, 4 scenes of ~15s each).
Each scene should include clear visual and audio cues.
Research Text: {text}
"""
        response = model.generate_content(prompt)
        if response and response.text:
            return response.text
        else:
            return None
    except Exception as e:
        st.error(f"Error generating 1min storyboard: {e}")
        return None

def generate_3min_storyboard(text):
    """
    Use the Gemini API to generate a 3-minute storyboard with 5-second interval breakdowns.
    """
    try:
        model = genai.GenerativeModel("gemini-pro")
        prompt = f"""
You are an expert in creating structured video storyboard scripts.
Generate a concise 3-minute storyboard (~180 seconds total) with 5-second interval breakdowns.
Include 6-7 scenes with clear visual and audio cues and an impactful call to action.
Research Text: {text}
"""
        response = model.generate_content(prompt)
        if response and response.text:
            return response.text
        else:
            return None
    except Exception as e:
        st.error(f"Error generating 3min storyboard: {e}")
        return None

def parse_storyboard(text):
    """
    Parse the generated storyboard into individual scenes.
    """
    scenes = []
    scene_blocks = re.split(r"(Scene\s+\d+)", text)
    for i in range(1, len(scene_blocks), 2):
        scene_title = scene_blocks[i].strip()
        details = scene_blocks[i+1].strip() if i+1 < len(scene_blocks) else ""
        lines = details.splitlines()
        visual = lines[0] if lines else "Visual details not provided."
        audio = lines[1] if len(lines) > 1 else "Audio details not provided."
        scenes.append({"scene": scene_title, "visual": visual, "audio": audio})
    return scenes

def create_full_video(parsed_scenes):
    """
    Generate a final vertical video from parsed storyboard scenes.
    This function fetches a short video clip from Pexels based on visual cues.
    """
    PEXELS_API_KEY = os.getenv("PEXELS_API_KEY", "YOUR PEXELS API KEY")  # Set your Pexels API key as env var
    def fetch_pexels_video(query):
        url = "https://api.pexels.com/videos/search"
        headers = {"Authorization": PEXELS_API_KEY}
        params = {"query": query, "per_page": 1}
        try:
            response = requests.get(url, headers=headers, params=params)
            response.raise_for_status()
            data = response.json()
            if data.get("videos"):
                return data["videos"][0]["video_files"][0]["link"]
        except Exception as e:
            st.warning(f"Error fetching video for '{query}': {e}")
        return None

    def make_vertical(video_clip):
        target_size = (1080, 1920)
        if video_clip.size[0] > video_clip.size[1]:
            new_width = int(video_clip.size[1] * 9 / 16)
            video_clip = video_clip.crop(x_center=video_clip.size[0] // 2, width=new_width)
        else:
            video_clip = video_clip.resize(height=1920)
        return video_clip.set_position("center").resize(target_size)

    video_clips = []
    for idx, scene in enumerate(parsed_scenes):
        st.write(f"Processing {scene['scene']}...")
        audio_filename = f"audio_{idx}.mp3"
        try:
            tts_audio = gTTS(scene["audio"], lang='en')
            tts_audio.save(audio_filename)
            audio_clip = AudioFileClip(audio_filename)
            scene_duration = audio_clip.duration
        except Exception as e:
            st.warning(f"Error generating audio for {scene['scene']}: {e}")
            continue
        video_url = fetch_pexels_video(scene["visual"])
        if video_url:
            video_filename = f"video_{idx}.mp4"
            try:
                vid_response = requests.get(video_url)
                with open(video_filename, "wb") as f:
                    f.write(vid_response.content)
                video_clip = VideoFileClip(video_filename).subclip(0, min(scene_duration, 5))
                video_clip = make_vertical(video_clip)
            except Exception as e:
                st.warning(f"Error processing video for {scene['scene']}: {e}")
                video_clip = None
        else:
            video_clip = None
        if not video_clip:
            video_clip = TextClip(scene["visual"], fontsize=50, color='white', bg_color='black', size=(1080, 1920)).set_duration(scene_duration)
        video_clip = video_clip.set_audio(audio_clip)
        video_clips.append(video_clip)
        time.sleep(1)
    if not video_clips:
        st.error("No valid scenes processed.")
        return None
    final_video = concatenate_videoclips(video_clips, method="compose")
    output_filename = "final_vertical_story_video.mp4"
    final_video.write_videofile(output_filename, codec="libx264", audio_codec="aac")
    return output_filename

# -------------------------
# STREAMLIT USER INTERFACE
# -------------------------
st.title("Research Paper Processing Pipeline")
st.subheader("Upload your research paper to get started")

# File Uploader
uploaded_pdf = st.file_uploader("Upload your research paper (PDF)", type=["pdf"])

if uploaded_pdf is not None:
    st.success("PDF uploaded successfully!")
    # Read the file once
    file_bytes = uploaded_pdf.getvalue()
    # Write file content to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
        tmp_file.write(file_bytes)
        tmp_file_path = tmp_file.name

    st.write("Extracting text from PDF...")
    extracted_text = extract_text_without_references(tmp_file_path)
    os.remove(tmp_file_path)  # Cleanup temporary file

    if extracted_text:
        st.text_area("Extracted Text", extracted_text, height=200)
        
        st.header("Choose a Processing Option")
        task = st.radio("Select a task:", 
                        ("Bullet Points Summarization", "PPT Generation", "Podcast Production", "Video Creation"))
        
        if task == "Bullet Points Summarization":
            with st.spinner("Generating summary..."):
                summary, refs = summarize_research_paper(extracted_text)
            st.subheader("Bullet Points Summary")
            st.text_area("Summary", summary, height=300)
            if refs:
                st.write("References Detected:", ", ".join(refs))
        
        elif task == "PPT Generation":
            with st.spinner("Creating PowerPoint presentation..."):
                ppt_path = create_ppt_from_text(extracted_text)
            st.success("PPT generated!")
            with open(ppt_path, "rb") as f:
                st.download_button("Download PPT", f, file_name=ppt_path)
        
        elif task == "Podcast Production":
            with st.spinner("Generating podcast script..."):
                podcast_script, refs = process_research_for_podcast(extracted_text)
            st.subheader("Podcast Script")
            st.text_area("Script", podcast_script, height=300)
            if refs:
                st.write("References Detected:", ", ".join(refs))
            if st.checkbox("Display Transcript"):
                st.text_area("Podcast Transcript", podcast_script, height=300)
            with st.spinner("Converting script to audio..."):
                podcast_audio_path = generate_gtts_audio(podcast_script, output_path="podcast_audio.mp3")
            if podcast_audio_path and os.path.exists(podcast_audio_path):
                st.audio(podcast_audio_path, format="audio/mp3")
                with open(podcast_audio_path, "rb") as f:
                    st.download_button("Download Podcast Audio", f, file_name=podcast_audio_path)
        
        elif task == "Video Creation":
            duration_option = st.selectbox("Select video storyboard duration:", ("1-Minute Storyboard", "3-Minute Storyboard"))
            if duration_option == "1-Minute Storyboard":
                with st.spinner("Generating 1-minute storyboard..."):
                    storyboard_text = generate_1min_storyboard(extracted_text)
            else:
                with st.spinner("Generating 3-minute storyboard..."):
                    storyboard_text = generate_3min_storyboard(extracted_text)
            if storyboard_text:
                st.subheader("Generated Storyboard")
                st.text_area("Storyboard", storyboard_text, height=300)
                parsed_scenes = parse_storyboard(storyboard_text)
                if parsed_scenes:
                    with st.spinner("Generating video from storyboard..."):
                        video_file = create_full_video(parsed_scenes)
                    if video_file and os.path.exists(video_file):
                        st.video(video_file)
                        with open(video_file, "rb") as f:
                            st.download_button("Download Video", f, file_name=video_file)
                else:
                    st.error("Failed to parse storyboard into scenes.")
            else:
                st.error("Storyboard generation failed.")
    else:
        st.error("No text extracted. Please check the PDF file.")
else:
    st.info("Please upload a PDF research paper to begin processing.")
